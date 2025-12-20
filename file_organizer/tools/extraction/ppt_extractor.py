# from .base_extractor import BaseExtractor

# class PPTExtractor(BaseExtractor):
#     def extract(self, file_path: str) -> str:
#         """Extract text from PowerPoint files."""
#         try:
#             from pptx import Presentation
#             prs = Presentation(file_path)
#             return "\n".join(
#                 shape.text 
#                 for slide in prs.slides 
#                 for shape in slide.shapes 
#                 if hasattr(shape, "text") and shape.text
#             )
#         except ImportError:
#             raise ImportError("Install python-pptx: pip install python-pptx")



from .base_extractor import BaseExtractor
import logging

logger = logging.getLogger(__name__)

class PPTExtractor(BaseExtractor):
    def extract(self, file_path: str) -> dict:
        """
        Extract text from PowerPoint files.
        
        Returns:
            dict: Dictionary containing 'content', 'success', and 'error' keys
        """
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            
            # Extract text from all slides and shapes
            content_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content_parts.append(shape.text.strip())
            
            content = "\n\n".join(content_parts)
            
            return {
                'content': content,
                'success': True,
                'file_path': file_path
            }
            
        except ImportError as e:
            error_msg = "Install python-pptx: pip install python-pptx"
            logger.error(error_msg)
            return {
                'content': '',
                'success': False,
                'error': error_msg,
                'file_path': file_path
            }
        except Exception as e:
            error_msg = f"Error extracting text from PowerPoint {file_path}: {str(e)}"
            logger.error(error_msg)
            return {
                'content': '',
                'success': False,
                'error': error_msg,
                'file_path': file_path
            }